#!/usr/bin/env python3
"""
Extract all content from a PowerPoint file (.pptx).
Returns a JSON structure with slides, text, and images.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]

Requires: pip install python-pptx
"""

import json
import os
import sys
from pptx import Presentation
from pathlib import Path


def safe_user_path(path_value, base_dir="."):
    """Resolve a CLI path under the current workspace."""
    if base_dir != ".":
        raise ValueError("Custom base directories are not supported for CLI paths")
    base_path = Path.cwd().resolve()
    resolved_path = Path(path_value).expanduser().resolve()
    try:
        resolved_path.relative_to(base_path)
    except ValueError as exc:
        raise ValueError(f"Path escapes allowed directory: {path_value}") from exc
    return resolved_path


def extract_pptx(file_path, output_dir="."):
    """
    Extract all content from a PowerPoint file.
    Returns a list of slide data dicts with text, images, and notes.
    """
    prs = Presentation(file_path)
    slides_data = []

    # Create assets directory for extracted images
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "number": slide_num + 1,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        for shape in slide.shapes:
            # Extract text content
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data["title"] = shape.text
                else:
                    slide_data["content"].append(
                        {"type": "text", "content": shape.text}
                    )

            # Extract images
            if shape.shape_type == 13:  # Picture type
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                image_path = os.path.join(assets_dir, image_name)

                with safe_user_path(image_path).open("wb") as f:
                    f.write(image_bytes)

                slide_data["images"].append(
                    {
                        "path": f"assets/{image_name}",
                        "width": shape.width,
                        "height": shape.height,
                    }
                )

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data["notes"] = notes_frame.text

        slides_data.append(slide_data)

    return slides_data


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract-pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = safe_user_path(sys.argv[2]) if len(sys.argv) > 2 else "."

    slides = extract_pptx(input_file, output_dir)

    # Write extracted data as JSON
    output_path = os.path.join(output_dir, "extracted-slides.json")
    with safe_user_path(output_path).open("w") as f:
        f.write(json.dumps(slides, indent=2))

    print(f"Extracted {len(slides)} slides to {output_path}")
    for s in slides:
        img_count = len(s["images"])
        print(f"  Slide {s['number']}: {s['title'] or '(no title)'} — {img_count} image(s)")
